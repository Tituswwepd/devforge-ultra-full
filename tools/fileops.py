import os, csv, zipfile, io
from pathlib import Path
from typing import Optional, List
from PIL import Image

# Optional deps
try:
    from reportlab.platypus import SimpleDocTemplate, Paragraph
    from reportlab.lib.styles import getSampleStyleSheet
except Exception:
    SimpleDocTemplate = None

try:
    import pypandoc
except Exception:
    pypandoc = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None

BASE = Path(__file__).resolve().parents[1]
OUT  = BASE / "out"
OUT.mkdir(exist_ok=True)

def _safe_name(name: str, default: str) -> str:
    name = (name or default).strip().replace("\\","/").split("/")[-1]
    return name or default

def make_txt(text: str, fname: str = "out.txt") -> str:
    p = OUT / _safe_name(fname, "out.txt")
    p.write_text(text, encoding="utf-8")
    return str(p)

def make_pdf(text: str, fname: str = "out.pdf") -> str:
    p = OUT / _safe_name(fname, "out.pdf")
    if not SimpleDocTemplate:
        p.write_text(text, encoding="utf-8")
        return str(p)
    doc = SimpleDocTemplate(str(p))
    styles = getSampleStyleSheet()
    story = [Paragraph(text.replace("\n","<br/>"), styles["Normal"])]
    doc.build(story)
    return str(p)

def make_docx(text: str, fname: str = "out.docx") -> str:
    p = OUT / _safe_name(fname, "out.docx")
    if not Document:
        # fallback: just write .docx as plain text (not real docx)
        p.write_text(text, encoding="utf-8")
        return str(p)
    d = Document()
    for line in text.splitlines():
        d.add_paragraph(line if line.strip() else "")
    d.save(str(p))
    return str(p)

def make_csv(rows: List[List[str]], fname: str = "table.csv") -> str:
    p = OUT / _safe_name(fname, "table.csv")
    with p.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        for r in rows:
            w.writerow(r)
    return str(p)

def make_pptx(title: str, bullets: List[str], fname: str = "slides.pptx") -> str:
    p = OUT / _safe_name(fname, "slides.pptx")
    if not Presentation:
        # fallback: text file
        p = OUT / (Path(fname).stem + ".txt")
        p.write_text("TITLE: " + title + "\n" + "\n".join(f"- {b}" for b in bullets), encoding="utf-8")
        return str(p)
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "\n".join(bullets) if bullets else ""
    prs.save(str(p))
    return str(p)

def convert_file(src: str, fmt: str) -> str:
    src_path = Path(src)
    if not src_path.exists():
        raise FileNotFoundError(f"Source not found: {src}")
    out = OUT / (src_path.stem + f".{fmt}")
    if not pypandoc:
        if src_path.suffix.lower().lstrip(".") == fmt.lower():
            out.write_bytes(src_path.read_bytes())
            return str(out)
        raise RuntimeError("pypandoc not installed; cannot convert.")
    pypandoc.convert_file(str(src_path), fmt, outputfile=str(out), extra_args=["--standalone"])
    return str(out)

def zip_files(paths: List[str], fname: str = "bundle.zip") -> str:
    z = OUT / _safe_name(fname, "bundle.zip")
    with zipfile.ZipFile(z, "w", zipfile.ZIP_DEFLATED) as outz:
        for p in paths:
            pp = Path(p)
            if pp.exists():
                outz.write(pp, arcname=pp.name)
    return str(z)

def unzip_file(zip_path: str, dest: Optional[str] = None) -> str:
    zpath = Path(zip_path)
    if not zpath.exists():
        raise FileNotFoundError(f"Zip not found: {zip_path}")
    dest_dir = Path(dest) if dest else (OUT / zpath.stem)
    dest_dir.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(zpath, "r") as z:
        z.extractall(dest_dir)
    return str(dest_dir)

def list_out_files() -> list[dict]:
    items = []
    for p in OUT.glob("*"):
        if p.is_file():
            items.append({"name": p.name, "size": p.stat().st_size})
    return sorted(items, key=lambda x: x["name"].lower())
