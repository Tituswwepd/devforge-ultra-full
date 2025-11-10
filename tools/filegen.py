# tools/filegen.py
from pathlib import Path
from io import BytesIO
import zipfile
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

def _safe_name(name: str, default: str) -> str:
    name = (name or "").strip()
    return name if name else default

def make_txt(name: str, content: str, out_dir: Path) -> Path:
    fname = _safe_name(name, "out.txt")
    if not fname.lower().endswith(".txt"):
        fname += ".txt"
    out = out_dir / fname.replace("/", "_").replace("\\", "_")
    out.write_text(content or "", encoding="utf-8")
    return out

def make_pdf(name: str, content: str, out_dir: Path) -> Path:
    fname = _safe_name(name, "out.pdf")
    if not fname.lower().endswith(".pdf"):
        fname += ".pdf"
    out = out_dir / fname.replace("/", "_").replace("\\", "_")
    c = canvas.Canvas(str(out), pagesize=A4)
    width, height = A4
    textobj = c.beginText(50, height - 50)
    for line in (content or "").splitlines():
        textobj.textLine(line)
    c.drawText(textobj)
    c.showPage()
    c.save()
    return out

def make_docx(name: str, content: str, out_dir: Path) -> Path:
    fname = _safe_name(name, "out.docx")
    if not fname.lower().endswith(".docx"):
        fname += ".docx"
    out = out_dir / fname.replace("/", "_").replace("\\", "_")
    doc = Document()
    for line in (content or "").splitlines():
        doc.add_paragraph(line)
    doc.save(out)
    return out

def make_csv(name: str, content: str, out_dir: Path) -> Path:
    fname = _safe_name(name, "out.csv")
    if not fname.lower().endswith(".csv"):
        fname += ".csv"
    out = out_dir / fname.replace("/", "_").replace("\\", "_")
    out.write_text(content or "", encoding="utf-8")
    return out

def make_pptx(name: str, content: str, out_dir: Path) -> Path:
    fname = _safe_name(name, "out.pptx")
    if not fname.lower().endswith(".pptx"):
        fname += ".pptx"
    out = out_dir / fname.replace("/", "_").replace("\\", "_")
    prs = Presentation()
    try:
        payload = content and content.strip()
        # Expect JSON: {"title": "...", "bullets": ["...","..."]}
        import json
        j = json.loads(payload) if payload else {}
        title = j.get("title", "Slides")
        bullets = j.get("bullets", ["Point one", "Point two"])
    except Exception:
        title = "Slides"
        bullets = (content or "Point one\nPoint two").splitlines()

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    if bullets:
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 1
    prs.save(out)
    return out

def make_zip(name: str, files: dict, out_dir: Path) -> Path:
    fname = _safe_name(name, "bundle.zip")
    if not fname.lower().endswith(".zip"):
        fname += ".zip"
    out = out_dir / fname.replace("/", "_").replace("\\", "_")
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for path, content in (files or {}).items():
            z.writestr(path, content if isinstance(content, str) else str(content))
    return out
